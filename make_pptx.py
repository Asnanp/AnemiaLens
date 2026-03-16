from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette (matches AnemiaLens app) ─────────────────────────────────────────
VOID    = RGBColor(0x02, 0x02, 0x08)
DEEP    = RGBColor(0x06, 0x06, 0x12)
CARD    = RGBColor(0x0A, 0x0A, 0x18)
CARD2   = RGBColor(0x10, 0x10, 0x22)
CRIMSON = RGBColor(0xC8, 0x00, 0x1E)
BRIGHT  = RGBColor(0xE8, 0x29, 0x4A)
WHITE   = RGBColor(0xF2, 0xF0, 0xEC)
MUTED   = RGBColor(0x8A, 0x88, 0x84)
DIM     = RGBColor(0x2A, 0x28, 0x30)
DIM2    = RGBColor(0x1E, 0x1C, 0x26)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
BLUE    = RGBColor(0x38, 0xBD, 0xF8)

W = prs.slide_width
H = prs.slide_height
blank = prs.slide_layouts[6]

# ── Core helpers ──────────────────────────────────────────────────────────────
def add_bg(slide, color=VOID):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def rect(slide, x, y, w, h, fill, lc=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb = lc; s.line.width = lw
    else:  s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=Pt(12), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def mono(slide, text, x, y, w, h, size=Pt(9), color=CRIMSON, bold=True, align=PP_ALIGN.LEFT):
    return txt(slide, text.upper(), x, y, w, h, size=size, bold=bold,
               color=color, align=align, font="Courier New")

def label(slide, text, x, y, w=Inches(6)):
    mono(slide, text, x, y, w, Inches(0.28), size=Pt(8), color=CRIMSON)

# Left crimson spine (used on every slide)
def spine(slide, w=Inches(0.14)):
    rect(slide, 0, 0, w, H, CRIMSON)

# Thin horizontal rule
def rule(slide, x, y, w, color=DIM, h=Inches(0.025)):
    rect(slide, x, y, w, h, color)

# Large ghost number (background watermark)
def ghost_num(slide, text, x, y, size=Pt(200)):
    txt(slide, text, x, y, Inches(5), Inches(3.5), size=size, bold=True,
        color=RGBColor(0x14, 0x12, 0x1E), align=PP_ALIGN.LEFT, font="Calibri")

# Stat chip: big value + small label
def stat_chip(slide, val, lbl, x, y, w=Inches(2.4), h=Inches(1.1),
              val_color=CRIMSON, border=DIM):
    rect(slide, x, y, w, h, CARD2, lc=border)
    txt(slide, val, x+Inches(0.18), y+Inches(0.08), w-Inches(0.3), Inches(0.55),
        size=Pt(28), bold=True, color=val_color, font="Calibri")
    mono(slide, lbl, x+Inches(0.18), y+Inches(0.65), w-Inches(0.3), Inches(0.3),
         size=Pt(7.5), color=MUTED, bold=False)

# Card with top crimson bar
def card(slide, x, y, w, h, title, body, accent=CRIMSON, title_size=Pt(13)):
    rect(slide, x, y, w, h, CARD2, lc=DIM)
    rect(slide, x, y, w, Inches(0.055), accent)
    txt(slide, title, x+Inches(0.22), y+Inches(0.14), w-Inches(0.4), Inches(0.42),
        size=title_size, bold=True, color=WHITE, font="Calibri")
    txt(slide, body, x+Inches(0.22), y+Inches(0.58), w-Inches(0.4), h-Inches(0.72),
        size=Pt(10), color=MUTED, font="Calibri")

# Left-accent card
def lcard(slide, x, y, w, h, title, body, accent=CRIMSON, title_size=Pt(13)):
    rect(slide, x, y, w, h, CARD2, lc=DIM)
    rect(slide, x, y, Inches(0.07), h, accent)
    txt(slide, title, x+Inches(0.22), y+Inches(0.14), w-Inches(0.36), Inches(0.42),
        size=title_size, bold=True, color=WHITE, font="Calibri")
    txt(slide, body, x+Inches(0.22), y+Inches(0.58), w-Inches(0.36), h-Inches(0.72),
        size=Pt(10), color=MUTED, font="Calibri")

# Diagonal accent line (decorative)
def diag_line(slide, x1, y1, x2, y2, color=DIM, w=Pt(0.5)):
    from pptx.util import Emu
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color; line.line.width = w

# Step node: circle-ish number badge + title + desc
def step_node(slide, num, title, desc, x, y, w=Inches(2.3), h=Inches(3.6), accent=CRIMSON):
    rect(slide, x, y, w, h, CARD2, lc=DIM)
    rect(slide, x, y, w, Inches(0.055), accent)
    # Number badge
    rect(slide, x+Inches(0.18), y+Inches(0.18), Inches(0.55), Inches(0.55), accent)
    txt(slide, num, x+Inches(0.18), y+Inches(0.16), Inches(0.55), Inches(0.55),
        size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    txt(slide, title, x+Inches(0.18), y+Inches(0.85), w-Inches(0.36), Inches(0.45),
        size=Pt(13), bold=True, color=WHITE, font="Calibri")
    txt(slide, desc, x+Inches(0.18), y+Inches(1.35), w-Inches(0.36), h-Inches(1.5),
        size=Pt(10), color=MUTED, font="Calibri")

# Arrow connector between steps
def arrow(slide, x, y):
    rect(slide, x, y+Inches(0.22), Inches(0.22), Inches(0.025), DIM)
    # arrowhead triangle approximation
    rect(slide, x+Inches(0.18), y+Inches(0.14), Inches(0.08), Inches(0.18), DIM)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)

# Ghost large "A" watermark top-right
ghost_num(s, "A", Inches(9.2), Inches(-0.8), size=Pt(320))

# Top rule
rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)

# Eyebrow
mono(s, "Hackathon 2026  ·  AI for Health  ·  Non-Invasive Screening",
     Inches(0.45), Inches(0.65), Inches(10), Inches(0.3), size=Pt(8), color=CRIMSON)

# Title — massive
txt(s, "AnemiaLens", Inches(0.45), Inches(1.3), Inches(10), Inches(2.0),
    size=Pt(96), bold=True, color=WHITE, font="Calibri")

# Crimson underline accent
rect(s, Inches(0.45), Inches(3.25), Inches(3.8), Inches(0.07), CRIMSON)

# Tagline
txt(s, "Smartphone-first anemia screening\npowered by computer vision + grounded GenAI.",
    Inches(0.45), Inches(3.45), Inches(7.8), Inches(1.1),
    size=Pt(17), color=MUTED, font="Calibri")

# Stat chips row — 4 chips
chips = [
    ("1.6B+", "People Affected", CRIMSON),
    ("92%",   "Sensitivity",     WHITE),
    ("$0",    "Marginal Cost",   GREEN),
    ("<60s",  "Time to Result",  AMBER),
]
for i, (val, lbl, vc) in enumerate(chips):
    stat_chip(s, val, lbl, Inches(0.45) + i * Inches(2.55), Inches(4.75),
              w=Inches(2.35), h=Inches(1.05), val_color=vc)

# Bottom rule + team tag
rule(s, Inches(0.35), Inches(6.85), Inches(12.6), DIM)
mono(s, "3-Minute Demo  ·  2026", Inches(9.5), Inches(6.95),
     Inches(3.5), Inches(0.3), size=Pt(8), color=DIM, bold=False)
mono(s, "Zero blood draw  ·  Works offline  ·  Any smartphone",
     Inches(0.45), Inches(6.95), Inches(7), Inches(0.3), size=Pt(8), color=DIM, bold=False)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)
ghost_num(s, "25%", Inches(7.5), Inches(0.5), size=Pt(180))

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)
label(s, "01  /  The Problem", Inches(0.45), Inches(0.65))
txt(s, "Bridging the Screening Gap.", Inches(0.45), Inches(0.95),
    Inches(9), Inches(1.1), size=Pt(46), bold=True, color=WHITE, font="Calibri")
rule(s, Inches(0.45), Inches(2.05), Inches(5.5), CRIMSON, h=Inches(0.04))

# Left: big stat block
rect(s, Inches(0.45), Inches(2.25), Inches(4.2), Inches(4.2), CARD2, lc=CRIMSON, lw=Pt(1.5))
rect(s, Inches(0.45), Inches(2.25), Inches(4.2), Inches(0.07), CRIMSON)
txt(s, "25%", Inches(0.6), Inches(2.4), Inches(3.9), Inches(1.5),
    size=Pt(88), bold=True, color=CRIMSON, font="Calibri")
txt(s, "of the global population\nhas anemia", Inches(0.6), Inches(3.9),
    Inches(3.9), Inches(0.65), size=Pt(14), color=WHITE, font="Calibri")
rule(s, Inches(0.6), Inches(4.65), Inches(3.8), DIM)
txt(s, "Yet most cases go undetected due to\ncost, distance, and systemic neglect.",
    Inches(0.6), Inches(4.75), Inches(3.8), Inches(0.9),
    size=Pt(10.5), color=MUTED, font="Calibri")

# Right: 4 pain point cards
pain = [
    ("Cost Barrier",    "Blood tests cost $5–$30 — unaffordable in low-income settings.", CRIMSON),
    ("Distance",        "Nearest clinic can be hours away for rural communities.",         AMBER),
    ("Delayed Results", "Results take days. Anemia worsens silently in the meantime.",    AMBER),
    ("Low Awareness",   "Symptoms — fatigue, dizziness — are vague and easy to dismiss.", MUTED),
]
for i, (title, desc, ac) in enumerate(pain):
    py = Inches(2.25) + i * Inches(1.06)
    lcard(s, Inches(4.9), py, Inches(8.1), Inches(0.95), title, desc, accent=ac, title_size=Pt(12))

rule(s, Inches(0.35), Inches(6.85), Inches(12.6), DIM)
mono(s, "Source: WHO Global Anaemia Estimates 2023", Inches(0.45), Inches(6.95),
     Inches(6), Inches(0.3), size=Pt(7.5), color=DIM, bold=False)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)
label(s, "02  /  The Solution", Inches(0.45), Inches(0.65))
txt(s, "Your phone is the lab.", Inches(0.45), Inches(0.95),
    Inches(10), Inches(1.1), size=Pt(46), bold=True, color=WHITE, font="Calibri")
rule(s, Inches(0.45), Inches(2.05), Inches(4.2), CRIMSON, h=Inches(0.04))

txt(s, "AnemiaLens analyzes the inner lower eyelid (conjunctiva). When hemoglobin drops,\n"
       "this tissue loses its pink color. Our model detects that pallor — non-invasively.",
    Inches(0.45), Inches(2.2), Inches(12.4), Inches(0.85),
    size=Pt(12.5), color=MUTED, font="Calibri")

# 3 pillar cards — equal width, full height
pillars = [
    ("Capture",  "Photograph the inner lower eyelid with any smartphone. No special hardware, no training required.",
     CRIMSON),
    ("Analyze",  "EfficientNet-B0 estimates hemoglobin (g/dL) from conjunctival pallor. Validated on 710 clinical specimens. Runs in <2s on CPU.",
     AMBER),
    ("Act",      "Structured clinical brief generated instantly. Triage band assigned. Share with a doctor in one tap.",
     GREEN),
]
for i, (title, desc, ac) in enumerate(pillars):
    px = Inches(0.45) + i * Inches(4.25)
    rect(s, px, Inches(3.2), Inches(4.0), Inches(3.5), CARD2, lc=ac, lw=Pt(1.2))
    rect(s, px, Inches(3.2), Inches(4.0), Inches(0.07), ac)
    # Large step number ghost
    txt(s, str(i+1), px+Inches(2.8), Inches(3.1), Inches(1.1), Inches(1.1),
        size=Pt(72), bold=True, color=RGBColor(0x16, 0x14, 0x24), font="Calibri")
    txt(s, title, px+Inches(0.22), Inches(3.42), Inches(3.6), Inches(0.5),
        size=Pt(20), bold=True, color=WHITE, font="Calibri")
    txt(s, desc, px+Inches(0.22), Inches(4.0), Inches(3.6), Inches(2.4),
        size=Pt(10.5), color=MUTED, font="Calibri")

# Bottom bar
rect(s, Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.5), CARD2, lc=DIM)
mono(s, "Zero blood draw  ·  Zero hardware  ·  Zero marginal cost  ·  Works offline  ·  Any smartphone",
     Inches(0.55), Inches(6.92), Inches(12.0), Inches(0.35),
     size=Pt(9), color=CRIMSON, bold=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)
label(s, "03  /  How It Works", Inches(0.45), Inches(0.65))
txt(s, "Five steps to screening clarity.", Inches(0.45), Inches(0.95),
    Inches(10), Inches(1.0), size=Pt(42), bold=True, color=WHITE, font="Calibri")
rule(s, Inches(0.45), Inches(1.95), Inches(5.0), CRIMSON, h=Inches(0.04))

steps = [
    ("01", "Capture",       "Photograph inner lower eyelid in bright daylight."),
    ("02", "Quality Gate",  "AI rejects blurry or misframed images before analysis."),
    ("03", "Vision AI",     "EfficientNet-B0 estimates hemoglobin from pallor."),
    ("04", "Symptom Fusion","Patient symptoms fuse with image biomarkers."),
    ("05", "Clinical Brief","Triage band + handoff summary generated instantly."),
]
for i, (num, title, desc) in enumerate(steps):
    sx = Inches(0.45) + i * Inches(2.52)
    step_node(s, num, title, desc, sx, Inches(2.15), w=Inches(2.3), h=Inches(4.0))
    # connector arrow between steps
    if i < 4:
        cx = sx + Inches(2.3)
        rect(s, cx, Inches(3.95), Inches(0.22), Inches(0.025), DIM)

rule(s, Inches(0.35), Inches(6.85), Inches(12.6), DIM)
mono(s, "Total time from capture to result: under 60 seconds",
     Inches(0.45), Inches(6.95), Inches(8), Inches(0.3), size=Pt(8), color=MUTED, bold=False)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)
label(s, "04  /  Technology", Inches(0.45), Inches(0.65))
txt(s, "The Intelligence Framework.", Inches(0.45), Inches(0.95),
    Inches(10), Inches(1.0), size=Pt(42), bold=True, color=WHITE, font="Calibri")
rule(s, Inches(0.45), Inches(1.95), Inches(4.8), CRIMSON, h=Inches(0.04))

layers = [
    ("Vision Layer",  "EfficientNet-B0",
     "Trained on 710 conjunctival specimens. Predicts hemoglobin (g/dL) and anemia risk score from a single image. Runs in <2s on CPU.",
     CRIMSON),
    ("GenAI Layer",   "Qwen-2.5  (Grounded)",
     "Generates safe, personalized clinical guidance. Constrained by deterministic medical rules — cannot hallucinate a diagnosis. Every output is rule-validated.",
     AMBER),
    ("Safety Layer",  "4-Band Triage System",
     "Low / Moderate / High Concern / Retake. Designed to prioritize patient safety over false confidence. Non-diagnostic language throughout.",
     GREEN),
]
for i, (layer, model, desc, color) in enumerate(layers):
    ly = Inches(2.15) + i * Inches(1.52)
    rect(s, Inches(0.45), ly, Inches(12.4), Inches(1.38), CARD2, lc=color, lw=Pt(1.2))
    rect(s, Inches(0.45), ly, Inches(0.1), Inches(1.38), color)
    mono(s, layer, Inches(0.7), ly+Inches(0.1), Inches(2.8), Inches(0.28),
         size=Pt(7.5), color=color, bold=True)
    txt(s, model, Inches(0.7), ly+Inches(0.36), Inches(3.8), Inches(0.55),
        size=Pt(20), bold=True, color=WHITE, font="Calibri")
    # Vertical divider
    rect(s, Inches(4.7), ly+Inches(0.15), Inches(0.025), Inches(1.08), DIM)
    txt(s, desc, Inches(4.9), ly+Inches(0.18), Inches(7.7), Inches(1.0),
        size=Pt(11), color=MUTED, font="Calibri")

# Stack tags
tags = ["React + Vite", "FastAPI", "EfficientNet-B0", "Qwen-2.5", "Python 3.11", "No cloud dependency"]
for i, tag in enumerate(tags):
    tx = Inches(0.45) + i * Inches(2.12)
    rect(s, tx, Inches(6.75), Inches(2.0), Inches(0.42), CARD2, lc=DIM)
    mono(s, tag, tx+Inches(0.1), Inches(6.78), Inches(1.8), Inches(0.32),
         size=Pt(8), color=MUTED, bold=False)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE DEMO OUTPUT  (data viz: Hb gauge + terminal)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)
label(s, "05  /  Live Demo Output", Inches(0.45), Inches(0.65))
txt(s, "Real screening result.", Inches(0.45), Inches(0.95),
    Inches(10), Inches(1.0), size=Pt(42), bold=True, color=WHITE, font="Calibri")
rule(s, Inches(0.45), Inches(1.95), Inches(3.8), CRIMSON, h=Inches(0.04))

# ── Left panel: triage result card ───────────────────────────────────────────
rect(s, Inches(0.45), Inches(2.15), Inches(6.0), Inches(4.55), CARD2, lc=AMBER, lw=Pt(1.5))
rect(s, Inches(0.45), Inches(2.15), Inches(6.0), Inches(0.07), AMBER)

mono(s, "Triage Result  ·  Moderate Risk", Inches(0.65), Inches(2.28),
     Inches(5.5), Inches(0.28), size=Pt(8), color=AMBER, bold=True)

# Hb value — huge
txt(s, "9.8", Inches(0.65), Inches(2.6), Inches(3.0), Inches(1.6),
    size=Pt(96), bold=True, color=AMBER, font="Calibri")
mono(s, "g/dL  Hemoglobin", Inches(0.65), Inches(4.2),
     Inches(5.5), Inches(0.3), size=Pt(9), color=MUTED, bold=False)

# Hb scale bar (visual gauge)
rule(s, Inches(0.65), Inches(4.6), Inches(5.5), DIM2, h=Inches(0.18))
# Fill: 9.8 out of 18 g/dL = ~54%
rect(s, Inches(0.65), Inches(4.6), Inches(2.97), Inches(0.18), AMBER)
mono(s, "0", Inches(0.65), Inches(4.82), Inches(0.4), Inches(0.22), size=Pt(7), color=DIM, bold=False)
mono(s, "9.8", Inches(3.4), Inches(4.82), Inches(0.6), Inches(0.22), size=Pt(7), color=AMBER, bold=True)
mono(s, "18 g/dL", Inches(5.9), Inches(4.82), Inches(0.7), Inches(0.22), size=Pt(7), color=DIM, bold=False)

txt(s, "Hemoglobin below normal range.\nClinical follow-up recommended within 1–2 weeks.",
    Inches(0.65), Inches(5.1), Inches(5.5), Inches(0.8),
    size=Pt(11), color=MUTED, font="Calibri")

# 3 mini stat chips below
mini_stats = [("68%", "Anemia Risk", CRIMSON), ("88%", "Confidence", WHITE), ("Band 2", "Triage", AMBER)]
for i, (val, lbl, vc) in enumerate(mini_stats):
    stat_chip(s, val, lbl, Inches(0.65) + i * Inches(1.9), Inches(6.0),
              w=Inches(1.75), h=Inches(0.85), val_color=vc, border=DIM)

# ── Right panel: terminal brief ───────────────────────────────────────────────
rect(s, Inches(6.7), Inches(2.15), Inches(6.3), Inches(4.55),
     RGBColor(0x00, 0x00, 0x00), lc=CRIMSON, lw=Pt(1.2))
# Terminal header bar
rect(s, Inches(6.7), Inches(2.15), Inches(6.3), Inches(0.38), CARD2)
mono(s, "CLINICAL HANDOFF BRIEF", Inches(6.9), Inches(2.22),
     Inches(5.8), Inches(0.28), size=Pt(8), color=CRIMSON, bold=True)
# Traffic light dots
for di, dc in enumerate([CRIMSON, AMBER, GREEN]):
    rect(s, Inches(12.3) - di * Inches(0.28), Inches(2.24), Inches(0.16), Inches(0.16), dc)

terminal_lines = [
    ("Provider:",    "EfficientNet-B0"),
    ("Risk Score:",  "68.4%  (Moderate)"),
    ("Hb Estimate:", "9.8 g/dL"),
    ("Confidence:",  "88%"),
    ("Triage Band:", "Band 2  —  Monitor"),
    ("Symptoms:",    "Fatigue, Dizziness"),
    ("Next Step:",   "CBC within 1–2 weeks"),
    ("Disclaimer:",  "Screening only. Not diagnostic."),
]
for i, (key, val) in enumerate(terminal_lines):
    ty = Inches(2.65) + i * Inches(0.44)
    mono(s, key, Inches(6.9), ty, Inches(1.9), Inches(0.38), size=Pt(10), color=CRIMSON, bold=True)
    mono(s, val, Inches(8.9), ty, Inches(3.9), Inches(0.38), size=Pt(10), color=WHITE, bold=False)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — IMPACT & SDGs
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)
ghost_num(s, "1.6B", Inches(7.0), Inches(0.2), size=Pt(160))

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)
label(s, "06  /  Impact", Inches(0.45), Inches(0.65))
txt(s, "Designed for the last mile.", Inches(0.45), Inches(0.95),
    Inches(10), Inches(1.0), size=Pt(42), bold=True, color=WHITE, font="Calibri")
rule(s, Inches(0.45), Inches(1.95), Inches(4.2), CRIMSON, h=Inches(0.04))

impacts = [
    ("Community Health Workers",
     "Deploy in rural India, sub-Saharan Africa, or any low-resource setting. No training beyond a 5-minute guide.",
     CRIMSON),
    ("Pregnant Women & Children",
     "The highest-risk groups — least likely to access a clinic. AnemiaLens meets them where they are.",
     BRIGHT),
    ("School & Field Nurses",
     "Rapid triage in schools, refugee camps, and disaster zones. Results in under 60 seconds.",
     AMBER),
    ("Telemedicine Platforms",
     "Structured clinical brief integrates directly into existing telehealth workflows via share/export.",
     GREEN),
]
for i, (title, desc, ac) in enumerate(impacts):
    ix = Inches(0.45) + (i % 2) * Inches(6.3)
    iy = Inches(2.15) + (i // 2) * Inches(2.05)
    lcard(s, ix, iy, Inches(6.0), Inches(1.85), title, desc, accent=ac, title_size=Pt(13))

# SDG bar
rect(s, Inches(0.45), Inches(6.38), Inches(12.4), Inches(0.52), CARD2, lc=DIM)
mono(s, "UN SDG 3: Good Health & Well-Being   ·   UN SDG 10: Reduced Inequalities   ·   UN SDG 1: No Poverty",
     Inches(0.55), Inches(6.46), Inches(12.0), Inches(0.35),
     size=Pt(9), color=CRIMSON, bold=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLOSE / CALL TO ACTION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
spine(s)

# Full-width crimson top bar
rect(s, 0, 0, W, Inches(0.12), CRIMSON)

# Ghost watermark
ghost_num(s, "AL", Inches(7.5), Inches(0.8), size=Pt(260))

rule(s, Inches(0.35), Inches(0.55), Inches(12.6), DIM)

txt(s, "AnemiaLens doesn't replace doctors.", Inches(0.45), Inches(1.1),
    Inches(12.4), Inches(1.2), size=Pt(38), bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Calibri")

txt(s, "It gets patients to doctors sooner,\nwith better information, at zero cost.",
    Inches(0.45), Inches(2.35), Inches(12.4), Inches(1.3),
    size=Pt(26), color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

rule(s, Inches(0.45), Inches(3.75), Inches(12.4), CRIMSON, h=Inches(0.05))

# Final 4 stats — large
final_stats = [
    ("1.6B+", "People we can reach",       CRIMSON),
    ("92%",   "Sensitivity",               WHITE),
    ("$0",    "Marginal cost / screening", GREEN),
    ("<60s",  "Time to result",            AMBER),
]
for i, (val, lbl, vc) in enumerate(final_stats):
    fx = Inches(0.45) + i * Inches(3.1)
    txt(s, val, fx, Inches(4.0), Inches(3.0), Inches(1.1),
        size=Pt(52), bold=True, color=vc, align=PP_ALIGN.CENTER, font="Calibri")
    mono(s, lbl, fx, Inches(5.1), Inches(3.0), Inches(0.35),
         size=Pt(8), color=MUTED, bold=False)

rule(s, Inches(0.45), Inches(5.6), Inches(12.4), DIM, h=Inches(0.025))

txt(s, "Thank you.", Inches(0.45), Inches(5.8),
    Inches(12.4), Inches(0.9), size=Pt(28), bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Calibri")

mono(s, "AnemiaLens  ·  Hackathon 2026  ·  AI for Health",
     Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.35),
     size=Pt(8), color=DIM, bold=False)

# ═════════════════════════════════════════════════════════════════════════════
prs.save("AnemiaLens_Presentation_v2.pptx")
print("Done  →  AnemiaLens_Presentation_v2.pptx")
